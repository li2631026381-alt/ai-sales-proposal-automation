# export_pptx.py
from pptx import Presentation

def export_to_pptx(data: dict, out_path: str) -> None:
    prs = Presentation()

    def add_title_slide(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_bullets(title: str, bullets: list[str]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b

    add_title_slide("AI Sales Proposal Automation", "Enterprise Edition – Demo Proposal")
    add_bullets("Executive Summary", [data.get("executive_summary", "")])
    add_bullets("Pain Point Analysis", data.get("pain_point_analysis", "").splitlines())

    sol = data.get("proposed_ai_solution", {})
    add_bullets("Proposed AI Solution", [
        f"Input: {sol.get('input','')}",
        f"Processing: {sol.get('processing','')}",
        f"Output: {sol.get('output','')}",
        "Modules: " + ", ".join(sol.get("modules", [])),
    ])

    plan = data.get("implementation_plan", [])
    plan_bullets = []
    for ph in plan:
        plan_bullets.append(ph["phase"])
        for it in ph["items"]:
            plan_bullets.append(f"  - {it}")
    add_bullets("Implementation Plan", plan_bullets)

    add_bullets("ROI Estimation", [data.get("roi_estimation_notes", "")])

    pricing = data.get("pricing_structure", [])
    pricing_bullets = []
    for t in pricing:
        pricing_bullets.append(f"{t['tier']} ({t['price_note']}): " + "; ".join(t["includes"]))
    add_bullets("Pricing Structure", pricing_bullets)

    risks = data.get("risks_mitigations", [])
    add_bullets("Risks & Mitigations", [f"{r['risk']} → {r['mitigation']}" for r in risks])
    add_bullets("Next Steps", data.get("next_steps", []))

    email = data.get("sales_email_draft", {}).get("en", "")
    add_bullets("Sales Email Draft (EN)", email.splitlines()[:10])

    prs.save(out_path)
